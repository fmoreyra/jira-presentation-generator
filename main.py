from jira_data_adapter import JiraDataAdapter
import os

from pptx import Presentation
from pptx.dml.color import RGBColor


from dotenv import load_dotenv
load_dotenv()

sprint = JiraDataAdapter(
    jira_url=os.environ['JIRA_URL'],
    jira_username=os.environ['JIRA_USERNAME'],
    jira_password=os.environ['JIRA_API_KEY']
).adapt()

sprint.print_tickets()

prs = Presentation('template_demo.pptx')
title_slide_layout = prs.slide_layouts[0]
epic_slide_layout = prs.slide_layouts[8]
ticket_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title

title.text = f"Demo - {sprint.name}"

EPIC_TITLE_INDEX = 0

TICKET_LAYOUT_INDEX = 2

TICKET_TITLE_INDEX = 0
EPIC_NAME_INDEX = 1
TICKET_ID_INDEX = 2
PERSON_IN_CHARGE_INDEX = 3
STATUS_INDEX = 4

epics = []
for ticket in sprint.tickets:
    if ticket.epic_name not in epics:
        epic_slide = prs.slides.add_slide(epic_slide_layout)
        epic_slide.shapes[EPIC_TITLE_INDEX].text = f"Épica: {ticket.epic_name}"
        epics.append(ticket.epic_name)
    ticket_slide = prs.slides.add_slide(prs.slide_layouts[TICKET_LAYOUT_INDEX])
    ticket_slide.shapes[TICKET_TITLE_INDEX].text = ticket.summary
    ticket_slide.shapes[EPIC_NAME_INDEX].text = f"Épica: {ticket.epic_name}"
    ticket_slide.shapes[TICKET_ID_INDEX].text = ticket.jira_id
    ticket_slide.shapes[PERSON_IN_CHARGE_INDEX].text = f"Responsable: {ticket.person_in_charge}"
    ticket_slide.shapes[STATUS_INDEX].text = f"Status: {ticket.status}"
    ticket_slide.shapes[STATUS_INDEX].fill.solid()
    ticket_slide.shapes[STATUS_INDEX].fill.fore_color.rgb = RGBColor(*ticket.status_color_rgb)

prs.save(f'Demo - {sprint.name}.pptx')
